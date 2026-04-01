#!/usr/bin/env python3
"""Build a Canva-friendly editable PPTX for thesis defense slides.

Usage:
    ./venv/bin/python build_editable_pptx.py
"""

from __future__ import annotations

from pathlib import Path
import shutil
import subprocess

from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def add_textbox(slide, left, top, width, height, text_lines, font_size=24, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    for idx, line in enumerate(text_lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
    return box


def add_title(slide, title):
    add_textbox(slide, 0.35, 0.1, 12.6, 0.7, [title], font_size=28, bold=True)


def add_picture_contain(slide, image_path: Path, left, top, width, height):
    box_w = Inches(width)
    box_h = Inches(height)
    box_l = Inches(left)
    box_t = Inches(top)

    with Image.open(image_path) as img:
        iw, ih = img.size

    img_aspect = iw / ih
    box_aspect = box_w / box_h

    if img_aspect > box_aspect:
        pic_w = box_w
        pic_h = int(box_w / img_aspect)
        pic_l = box_l
        pic_t = box_t + (box_h - pic_h) // 2
    else:
        pic_h = box_h
        pic_w = int(box_h * img_aspect)
        pic_t = box_t
        pic_l = box_l + (box_w - pic_w) // 2

    slide.shapes.add_picture(str(image_path), pic_l, pic_t, width=pic_w, height=pic_h)


def set_table_font(table, size_pt=15, header_bold=True):
    for r in range(len(table.rows)):
        for c in range(len(table.columns)):
            cell = table.cell(r, c)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(size_pt)
                    if r == 0 and header_bold:
                        run.font.bold = True


def set_table_alignment(table, first_col_align=PP_ALIGN.LEFT, other_align=PP_ALIGN.CENTER):
    for r in range(len(table.rows)):
        for c in range(len(table.columns)):
            align = first_col_align if c == 0 else other_align
            cell = table.cell(r, c)
            for p in cell.text_frame.paragraphs:
                p.alignment = align


def set_cell_bold(cell, bold=True):
    for p in cell.text_frame.paragraphs:
        for run in p.runs:
            run.font.bold = bool(bold)


def add_visual_placeholder(slide, left, top, width, height, label):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(246, 248, 252)
    shape.line.color.rgb = RGBColor(108, 117, 125)
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = label
    p1.alignment = PP_ALIGN.CENTER
    p1.runs[0].font.bold = True
    p1.runs[0].font.size = Pt(24)
    p2 = tf.add_paragraph()
    p2.text = "Replace with your custom diagram"
    p2.alignment = PP_ALIGN.CENTER
    p2.runs[0].font.size = Pt(16)


def ensure_elephants_slide_image(figures_dir: Path, presentation_dir: Path) -> Path:
    pdf_path = figures_dir / "elephants_train_SAFE.pdf"
    png_path = presentation_dir / "elephants_train_SAFE_slide.png"
    if png_path.exists():
        return png_path
    cmd = [
        "pdftoppm",
        "-png",
        "-singlefile",
        str(pdf_path),
        str(png_path.with_suffix("")),
    ]
    subprocess.run(cmd, check=True)
    return png_path


def build() -> Path:
    here = Path(__file__).resolve().parent
    figs = (here / "../Figures").resolve()
    out = here / "thesis_defense_editable.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1: title
    s = prs.slides.add_slide(blank)
    add_textbox(
        s,
        0.7,
        1.2,
        12.0,
        1.9,
        [
            "Deep Learning-Based Animal Re-Identification",
            "for Non-Invasive Wildlife Monitoring and Conservation",
        ],
        font_size=40,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s,
        0.7,
        3.7,
        12.0,
        1.1,
        ["MSc Thesis Defense (15 min)", "Matej Maric - February 2026"],
        font_size=22,
        align=PP_ALIGN.CENTER,
    )

    # Slide 2: motivation
    s = prs.slides.add_slide(blank)
    add_title(s, "Motivation: Why a Species-Agnostic Pipeline?")
    add_textbox(
        s,
        0.65,
        1.25,
        5.4,
        5.8,
        [
            "- Unified pipeline across species",
            "- Minimal supervision, practical runtime",
            "- Re-ID as a building block for counting",
        ],
        font_size=26,
    )
    add_picture_contain(s, figs / "dataset_examples_grid.JPG", 6.2, 1.05, 6.8, 6.25)

    # Slide 3: datasets overview
    s = prs.slides.add_slide(blank)
    add_title(s, "Datasets: WildlifeReID-10k (7 Used Here)")
    add_textbox(
        s,
        0.65,
        0.95,
        12.2,
        0.95,
        [
            "WildlifeReID-10k: 37 datasets, 33 species, 140k images, 10.7k individuals.",
            "Evaluated on 7 datasets spanning easy to hard cases (closed-set).",
        ],
        font_size=18,
    )
    table_shape = s.shapes.add_table(8, 4, Inches(0.55), Inches(1.9), Inches(12.25), Inches(4.9))
    table = table_shape.table
    headers = ["Dataset", "# images", "# individuals", "split type"]
    rows = [
        ["ATRW*", "5,415", "182", "Original MD split"],
        ["CZoo*", "4,662", "71", "Original MD split"],
        ["SealID*", "2,080", "57", "Original MD split"],
        ["CowDataset", "1,485", "13", "Time-aware split"],
        ["Chicks4FreeID", "1,146", "50", "Similarity-aware split"],
        ["SeaStarReID2023", "2,187", "95", "Time-aware split"],
        ["ELPephants", "2,078", "276", "Random split"],
    ]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
    set_table_font(table, size_pt=16)
    set_table_alignment(table, first_col_align=PP_ALIGN.LEFT, other_align=PP_ALIGN.CENTER)
    add_textbox(
        s,
        0.65,
        6.9,
        12.2,
        0.5,
        ["* MD-trained datasets: use the exact training split to avoid additional leakage."],
        font_size=14,
    )

    # Slide 4: leakage control
    s = prs.slides.add_slide(blank)
    add_title(s, "Evaluation Rigor: Leakage Control")
    add_picture_contain(s, figs / "atrw_data_leak.png", 0.35, 1.0, 6.35, 5.9)
    add_textbox(
        s,
        7.0,
        1.55,
        5.8,
        3.7,
        ["- Time-aware splits", "- Similarity-aware splits", "- MD-train split alignment (ATRW/CZoo/SealID*)"],
        font_size=24,
    )

    elephants_img = ensure_elephants_slide_image(figs, here)

    # Slide 5: classification funnel bullets
    s = prs.slides.add_slide(blank)
    add_title(s, "Classification Funnel: Retrieve -> Rerank -> Verify")
    add_textbox(
        s,
        0.75,
        1.25,
        12.0,
        5.6,
        [
            "- Tier-1: retrieve candidates (Global + Fisher) and take the union",
            "- Tier-2: calibrate and fuse scores into same-ID probability",
            "- Tier-3: geometric verification (keypoints + RANSAC) on the shortlist",
        ],
        font_size=28,
    )

    # Slide 6: full pipeline figure
    s = prs.slides.add_slide(blank)
    add_title(s, "Pipeline Overview (Full)")
    add_picture_contain(s, figs / "classification_pipeline.png", 0.35, 0.95, 12.6, 6.2)

    # Slide 7: tier 1 concept
    s = prs.slides.add_slide(blank)
    add_title(s, "Tier 1: Candidate Retrieval (Concept)")
    add_textbox(
        s,
        0.85,
        1.25,
        12.0,
        5.6,
        [
            "- Two fast signals: Global embeddings + Fisher vectors",
            "- Take top-K from each and form an ordered union (deduped)",
            "- Goal: high-recall shortlist for expensive stages",
            "- Background removal helps stabilize local evidence (optional)",
        ],
        font_size=26,
    )

    # Slide 8: tier 1 visual
    s = prs.slides.add_slide(blank)
    add_title(s, "Tier 1: Candidate Retrieval (Visualization)")
    add_picture_contain(s, figs / "tier1_union_flow.png", 0.35, 0.95, 12.6, 6.2)

    # Slide 9: tier 2 concept
    s = prs.slides.add_slide(blank)
    add_title(s, "Tier 2: Calibration and Fusion (Concept)")
    add_textbox(
        s,
        0.85,
        1.25,
        12.0,
        5.6,
        [
            "- Raw similarities are not comparable across modalities",
            "- Per-dataset calibration: score -> P(same | score)",
            "- Fuse calibrated Global + Fisher (simple average)",
            "- Output: Tier-2 reranking of the union shortlist",
        ],
        font_size=26,
    )

    # Slide 10: tier 2 visual
    s = prs.slides.add_slide(blank)
    add_title(s, "Tier 2: Calibration and Fusion (Visualization)")
    add_picture_contain(s, figs / "tier2_fusion_flow.png", 0.35, 0.95, 12.6, 6.2)

    # Slide 11: tier 3 concept
    s = prs.slides.add_slide(blank)
    add_title(s, "Tier 3: Geometric Verification (Concept)")
    add_textbox(
        s,
        0.85,
        1.25,
        12.0,
        5.6,
        [
            "- Run only on a small shortlist (most expensive stage)",
            "- Match local keypoints and estimate geometric consistency",
            "- Inlier count provides strong evidence against false matches",
            "- Especially helpful on hard datasets (e.g., ELPephants)",
        ],
        font_size=26,
    )

    # Slide 12: tier 3 visual
    s = prs.slides.add_slide(blank)
    add_title(s, "Tier 3: Geometric Verification (Visualization)")
    add_picture_contain(s, figs / "tier3_gv_flow.png", 0.35, 0.95, 12.6, 6.2)

    # Slide 13: classification results
    s = prs.slides.add_slide(blank)
    add_title(s, "Classification Results (Mean over 7 datasets)")
    table_shape = s.shapes.add_table(5, 4, Inches(0.65), Inches(1.55), Inches(7.2), Inches(3.1))
    table = table_shape.table
    headers = ["Metric", "Global", "WildFusion", "This thesis"]
    rows = [
        ["Top-1 (%)", "71.75", "88.18", "86.73"],
        ["Top-5 (%)", "79.58", "91.91", "90.70"],
        ["F1 (%)", "71.51", "87.35", "85.70"],
        ["Runtime (min)", "0.06", "214.52", "42.53"],
    ]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
    set_table_font(table, size_pt=16)
    set_table_alignment(table, first_col_align=PP_ALIGN.LEFT, other_align=PP_ALIGN.CENTER)
    add_textbox(
        s,
        8.1,
        1.45,
        4.9,
        4.5,
        [
            "- WildFusion: strong calibrated fusion baseline",
            "- GV is expensive -> applied only on a shortlist",
            "- Key point: near-WildFusion accuracy at ~5x lower runtime",
        ],
        font_size=24,
    )

    # Slide 14: selected datasets table
    s = prs.slides.add_slide(blank)
    add_title(s, "Classification Results (Selected Datasets)")
    cols_sel = ["Dataset", "Metric", "Global", "WildFusion", "Fisher", "G+F", "F+GV", "G+F+GV"]
    rows_sel = [
        ["ELPephants", "Top-1", "13.66", "49.31", "20.79", "21.58", "46.73", "54.46"],
        ["", "Top-5", "21.78", "57.03", "31.49", "32.48", "68.12", "64.95"],
        ["", "F1", "11.58", "45.19", "18.21", "19.26", "42.87", "49.71"],
        ["SealID*", "Top-1", "78.42", "97.60", "65.23", "79.38", "81.53", "85.13"],
        ["", "Top-5", "79.62", "98.56", "80.58", "82.73", "89.69", "88.97"],
        ["", "F1", "77.86", "97.67", "64.59", "78.47", "80.58", "84.40"],
        ["SeaStarReID2023", "Top-1", "47.91", "80.47", "68.84", "62.79", "77.21", "77.21"],
        ["", "Top-5", "73.49", "90.70", "85.58", "82.79", "85.58", "85.58"],
        ["", "F1", "45.71", "79.09", "67.92", "60.70", "76.40", "76.85"],
    ]
    table_shape = s.shapes.add_table(
        1 + len(rows_sel),
        len(cols_sel),
        Inches(0.25),
        Inches(1.35),
        Inches(12.85),
        Inches(5.2),
    )
    table = table_shape.table
    for c, h in enumerate(cols_sel):
        table.cell(0, c).text = h
    for r, row in enumerate(rows_sel, start=1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
    set_table_font(table, size_pt=12)
    set_table_alignment(table, first_col_align=PP_ALIGN.LEFT, other_align=PP_ALIGN.CENTER)
    # Metric column left-aligned.
    for r in range(0, 1 + len(rows_sel)):
        for p in table.cell(r, 1).text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
    # Bold header row + dataset column.
    for c in range(len(cols_sel)):
        set_cell_bold(table.cell(0, c), True)
    for r in range(1, 1 + len(rows_sel)):
        set_cell_bold(table.cell(r, 0), True)
    # Bold best per row (copied from thesis table).
    best_cells = [
        (1, 7),  # ELPephants Top-1: G+F+GV
        (2, 6),  # ELPephants Top-5: F+GV
        (3, 7),  # ELPephants F1:   G+F+GV
        (4, 3),  # SealID*: all best WildFusion
        (5, 3),
        (6, 3),
        (7, 3),  # SeaStar: all best WildFusion
        (8, 3),
        (9, 3),
    ]
    for (rr, cc) in best_cells:
        set_cell_bold(table.cell(rr, cc), True)
    add_textbox(
        s,
        0.65,
        6.65,
        12.2,
        0.6,
        ["GV gain (ELPephants Top-1): 21.58 (G+F) -> 54.46 (G+F+GV)."],
        font_size=16,
    )

    # Slide 15: ELPephants special case
    s = prs.slides.add_slide(blank)
    add_title(s, "ELPephants: A Stress-Test for Local + GV")
    add_textbox(
        s,
        0.65,
        1.2,
        6.3,
        4.9,
        ["- Smooth texture -> weak global cues", "- Strong pose/illumination/mud variation", "- Local evidence + GV is crucial", "- Best Top-1: 54.46% (vs 49.31% WildFusion)"],
        font_size=22,
    )
    add_picture_contain(s, elephants_img, 6.95, 1.1, 5.8, 5.85)

    # Slide 16: counting theory
    s = prs.slides.add_slide(blank)
    add_title(s, "Population Counting: HITL-NIS (Idea)")
    add_textbox(
        s,
        0.65,
        1.2,
        7.4,
        5.4,
        [
            "- Goal: estimate #individuals without full clustering",
            "- Model images as clique graph -> connected components",
            "- Query only a small number of pairs (oracle: same/different)",
            "- Use Nested Importance Sampling to reduce variance",
        ],
        font_size=24,
    )
    add_textbox(
        s,
        8.25,
        2.2,
        4.7,
        1.8,
        ["K = Σu 1 / (1 + d(u))"],
        font_size=30,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s,
        8.25,
        4.1,
        4.7,
        1.0,
        ["N vertices, M neighbors"],
        font_size=20,
        align=PP_ALIGN.CENTER,
    )

    # Slide 17: counting drawbacks
    s = prs.slides.add_slide(blank)
    add_title(s, "Counting Drawbacks: Human Error Matters")
    add_textbox(
        s,
        0.75,
        1.25,
        12.0,
        5.7,
        [
            "- Pair labels are highly imbalanced (mostly 'different')",
            "- False positives merge identities -> underestimation collapse",
            "- Confidence intervals capture sampling variance, not bias",
            "- Mitigation: positive confirmation (accept 'same' only after K votes)",
        ],
        font_size=26,
    )

    # Slide 18: counting results K=2
    s = prs.slides.add_slide(blank)
    add_title(s, "Population Results (K=2): Full Table")
    cols = ["", "ATRW", "CowDataset", "Chicks4FreeID", "CZoo", "ELPephants", "SealID", "SeaStarReID2023"]
    data = [
        ["#images", "5415", "1388", "1086", "2109", "2078", "2080", "2077"],
        ["GT", "182", "12", "48", "24", "274", "57", "91"],
        ["p=0.00", "278 [132, 425]", "12 [11, 13]", "50 [43, 58]", "24 [23, 26]", "334 [267, 401]", "56 [45, 68]", "106 [85, 128]"],
        ["0.02", "212 [124, 300]", "13 [11, 14]", "50 [42, 57]", "25 [23, 27]", "320 [259, 381]", "65 [47, 83]", "106 [87, 126]"],
        ["0.05", "178 [116, 240]", "13 [12, 14]", "47 [41, 54]", "26 [24, 27]", "178 [157, 199]", "53 [43, 62]", "93 [79, 106]"],
        ["0.10", "72 [58, 86]", "13 [12, 14]", "39 [34, 43]", "25 [23, 27]", "78 [73, 83]", "45 [37, 54]", "59 [52, 66]"],
        ["0.15", "39 [34, 45]", "12 [11, 14]", "27 [24, 30]", "22 [20, 24]", "40 [38, 43]", "33 [29, 38]", "37 [33, 41]"],
        ["0.30", "11 [10, 12]", "8 [7, 9]", "10 [9, 11]", "13 [11, 14]", "11 [10, 12]", "13 [11, 15]", "11 [10, 12]"],
    ]
    table_shape = s.shapes.add_table(1 + len(data), len(cols), Inches(0.25), Inches(1.35), Inches(12.85), Inches(5.55))
    table = table_shape.table
    for c, h in enumerate(cols):
        table.cell(0, c).text = h
    for r, row in enumerate(data, start=1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
    set_table_font(table, size_pt=10)
    set_table_alignment(table, first_col_align=PP_ALIGN.LEFT, other_align=PP_ALIGN.CENTER)
    # Bold headers + row labels
    for r in range(0, 1 + len(data)):
        set_cell_bold(table.cell(r, 0), True)
    # Bold cells where GT is inside CI (copied from thesis table formatting)
    bold_k2 = set()
    # p=0.00 and 0.02: all datasets bold
    for rr in [3, 4]:
        for cc in range(1, len(cols)):
            bold_k2.add((rr, cc))
    # p=0.05: all except ELPephants
    for cc in range(1, len(cols)):
        if cols[cc] != "ELPephants":
            bold_k2.add((5, cc))
    # p=0.10: CowDataset, CZoo
    for name in ["CowDataset", "CZoo"]:
        bold_k2.add((6, cols.index(name)))
    # p=0.15: CowDataset, CZoo
    for name in ["CowDataset", "CZoo"]:
        bold_k2.add((7, cols.index(name)))
    for (rr, cc) in bold_k2:
        set_cell_bold(table.cell(rr, cc), True)
    add_textbox(s, 0.65, 6.95, 12.2, 0.45, ["Bold: GT inside the 95% confidence interval."], font_size=14)

    # Slide 19: counting results K=1
    s = prs.slides.add_slide(blank)
    add_title(s, "Population Results (K=1): Full Table")
    data = [
        ["#images", "5415", "1388", "1086", "2109", "2078", "2080", "2077"],
        ["GT", "182", "12", "48", "24", "274", "57", "91"],
        ["p=0.00", "278 [132, 425]", "12 [11, 13]", "50 [43, 58]", "24 [23, 26]", "342 [298, 386]", "56 [45, 68]", "106 [85, 128]"],
        ["0.02", "40 [33, 46]", "10 [9, 11]", "25 [22, 28]", "19 [17, 20]", "42 [41, 44]", "29 [25, 34]", "36 [32, 40]"],
        ["0.05", "18 [16, 20]", "8 [7, 9]", "15 [13, 16]", "13 [12, 15]", "19 [18, 20]", "17 [15, 20]", "18 [16, 20]"],
        ["0.10", "10 [9, 11]", "6 [5, 7]", "9 [8, 10]", "9 [8, 10]", "10 [9, 10]", "11 [9, 12]", "10 [9, 11]"],
        ["0.15", "7 [6, 7]", "5 [4, 5]", "6 [5, 7]", "7 [6, 8]", "7 [6, 7]", "8 [7, 8]", "7 [6, 7]"],
        ["0.30", "3 [3, 4]", "3 [2, 4]", "3 [3, 4]", "4 [3, 5]", "3 [3, 4]", "4 [3, 5]", "3 [3, 4]"],
    ]
    table_shape = s.shapes.add_table(1 + len(data), len(cols), Inches(0.25), Inches(1.35), Inches(12.85), Inches(5.55))
    table = table_shape.table
    for c, h in enumerate(cols):
        table.cell(0, c).text = h
    for r, row in enumerate(data, start=1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
    set_table_font(table, size_pt=10)
    set_table_alignment(table, first_col_align=PP_ALIGN.LEFT, other_align=PP_ALIGN.CENTER)
    for r in range(0, 1 + len(data)):
        set_cell_bold(table.cell(r, 0), True)
    # Bold cells where GT is inside CI (copied from thesis table formatting)
    bold_k1 = set()
    # p=0.00: all except ELPephants
    for cc in range(1, len(cols)):
        if cols[cc] != "ELPephants":
            bold_k1.add((3, cc))
    for (rr, cc) in bold_k1:
        set_cell_bold(table.cell(rr, cc), True)
    add_textbox(s, 0.65, 6.95, 12.2, 0.45, ["Bold: GT inside the 95% confidence interval."], font_size=14)

    # Slide 20: conclusions
    s = prs.slides.add_slide(blank)
    add_title(s, "Conclusions")
    add_textbox(
        s,
        0.75,
        1.2,
        11.8,
        4.9,
        [
            "- Species-agnostic 3-tier Re-ID works across multiple benchmarks",
            "- Global + Fisher + GV gives strong accuracy/runtime trade-off",
            "- HITL-NIS is promising, but very sensitive to human error",
        ],
        font_size=28,
    )

    # Slide 21: future work
    s = prs.slides.add_slide(blank)
    add_title(s, "Future Work")
    add_textbox(
        s,
        0.75,
        1.2,
        11.8,
        4.9,
        [
            "- Validate counting with real annotators (not simulated error)",
            "- Better global backbone and stronger local matching for hard datasets",
            "- Link photographed-individual counts to true abundance/density",
        ],
        font_size=28,
    )

    # Slide 22: thank you
    s = prs.slides.add_slide(blank)
    add_textbox(
        s,
        0.75,
        0.75,
        11.8,
        1.3,
        ["Thank you", "Questions?"],
        font_size=44,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_picture_contain(s, figs / "segmentation_pipeline_vis.JPG", 1.5, 2.2, 10.3, 4.7)

    prs.save(out)
    shutil.copyfile(out, here / "thesis_defense_editable_canva.pptx")
    return out


def main():
    out = build()
    print(f"Wrote: {out}")


if __name__ == "__main__":
    main()
